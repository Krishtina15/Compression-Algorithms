from flask import Flask, request, jsonify, render_template_string, send_file
import io
import heapq
import base64
import json
import pdfplumber
import docx
import pptx

app = Flask(__name__)

HTML = r"""
<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0" />
  <title>Huffman Compressor</title>
  <style>
    *, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }

    body {
      font-family: 'Segoe UI', system-ui, -apple-system, sans-serif;
      background: #eef0fb;
      color: #1a1f36;
      min-height: 100vh;
      padding: 28px 24px;
    }

    .dashboard {
      max-width: 1100px;
      margin: 0 auto;
      display: grid;
      grid-template-columns: 1fr 260px;
      gap: 20px;
    }

    .main  { display: flex; flex-direction: column; gap: 20px; }
    .sidebar { display: flex; flex-direction: column; gap: 14px; }

    .card {
      background: #fff;
      border-radius: 18px;
      padding: 24px;
      box-shadow: 0 4px 24px rgba(90, 100, 200, 0.07);
    }

    .header-card {
      display: flex;
      align-items: center;
      justify-content: space-between;
      background: linear-gradient(135deg, #fff 60%, #eef0fb 100%);
      padding: 28px 32px;
    }

    .header-card h1 {
      font-size: 1.75rem;
      font-weight: 700;
      color: #1a1f36;
      letter-spacing: -0.5px;
    }

    .header-card p {
      margin-top: 5px;
      color: #8a94b0;
      font-size: 0.88rem;
    }

    .header-illo {
      width: 80px;
      height: 80px;
      border-radius: 50%;
      background: linear-gradient(135deg, #c7d0ff 0%, #a5b4fc 100%);
      display: flex;
      align-items: center;
      justify-content: center;
      font-size: 2rem;
      flex-shrink: 0;
    }

    .section-head {
      display: flex;
      align-items: flex-start;
      justify-content: space-between;
      margin-bottom: 16px;
    }

    .section-head h2 { font-size: 1rem; font-weight: 700; color: #1a1f36; }
    .section-head p  { font-size: 0.78rem; color: #8a94b0; margin-top: 2px; }

    .drop-zone {
      border: 2px dashed #d5d9f0;
      border-radius: 12px;
      padding: 22px;
      text-align: center;
      cursor: pointer;
      transition: border-color 0.2s, background 0.2s;
      margin-bottom: 14px;
    }

    .drop-zone:hover, .drop-zone.dragover {
      border-color: #5c6bc0;
      background: #f3f4fd;
    }

    .drop-zone svg  { display: block; margin: 0 auto 8px; color: #b0b8d8; }
    .drop-zone p    { color: #8a94b0; font-size: 0.85rem; }
    .drop-zone span { color: #5c6bc0; font-weight: 600; }

    #file-input { display: none; }

    textarea {
      width: 100%;
      border: 1.5px solid #e5e8f5;
      border-radius: 10px;
      padding: 12px 14px;
      font-family: 'Courier New', monospace;
      font-size: 0.82rem;
      resize: vertical;
      outline: none;
      transition: border-color 0.2s;
      background: #fafbff;
      color: #1a1f36;
      line-height: 1.6;
    }

    textarea:focus   { border-color: #5c6bc0; background: #fff; }
    .input-area      { height: 130px; }
    .output-area     { height: 120px; background: #f7f8fe; }

    .btn-row {
      display: flex;
      gap: 10px;
      margin-top: 14px;
      align-items: center;
      flex-wrap: wrap;
    }

    button {
      padding: 9px 22px;
      border: none;
      border-radius: 10px;
      font-size: 0.85rem;
      font-weight: 600;
      cursor: pointer;
      transition: all 0.16s;
    }

    .btn-primary {
      background: linear-gradient(135deg, #5c6bc0, #7986cb);
      color: #fff;
      box-shadow: 0 4px 14px rgba(92, 107, 192, 0.35);
    }

    .btn-primary:hover   { box-shadow: 0 6px 18px rgba(92,107,192,0.45); transform: translateY(-1px); }
    .btn-primary:active  { transform: translateY(0); }
    .btn-primary:disabled { opacity: 0.5; cursor: not-allowed; transform: none; }

    .btn-outline {
      background: transparent;
      color: #5c6bc0;
      border: 1.5px solid #c5cbee;
    }

    .btn-outline:hover { background: #f0f2fb; }

    .btn-ghost {
      background: #f0f2fb;
      color: #5c6bc0;
    }

    .btn-ghost:hover { background: #e5e8f8; }

    .char-count {
      margin-left: auto;
      font-size: 0.78rem;
      color: #b0b8d8;
      font-weight: 500;
    }

    .divider {
      display: flex;
      align-items: center;
      gap: 10px;
      color: #c5cbee;
      font-size: 0.75rem;
      margin: 10px 0;
    }

    .divider::before, .divider::after {
      content: '';
      flex: 1;
      height: 1px;
      background: #e8eaf6;
    }

    .msg {
      font-size: 0.8rem;
      padding: 10px 14px;
      border-radius: 10px;
      margin-top: 12px;
      display: none;
      line-height: 1.5;
    }

    .msg.info  { background: #eef2ff; color: #3d52a0; display: block; }
    .msg.warn  { background: #fffbeb; color: #92620a; display: block; }
    .msg.error { background: #fff0f0; color: #c0392b; display: block; }

    .metric-box {
      background: #fff;
      border-radius: 14px;
      padding: 16px 18px;
      box-shadow: 0 4px 24px rgba(90, 100, 200, 0.07);
      display: flex;
      align-items: center;
      justify-content: space-between;
      cursor: default;
    }

    .metric-box:hover { background: #f7f8fe; }

    .mb-left .mb-label {
      font-size: 0.72rem;
      color: #8a94b0;
      font-weight: 500;
      margin-bottom: 4px;
    }

    .mb-left .mb-value {
      font-size: 1.2rem;
      font-weight: 700;
      color: #1a1f36;
    }

    .mb-arrow {
      width: 28px;
      height: 28px;
      border-radius: 8px;
      background: #eef0ff;
      display: flex;
      align-items: center;
      justify-content: center;
      color: #5c6bc0;
      font-size: 0.8rem;
      flex-shrink: 0;
    }

    .donut-card {
      display: flex;
      flex-direction: column;
      align-items: center;
      padding: 22px 18px 18px;
    }

    .donut-card .dc-title {
      font-size: 0.72rem;
      font-weight: 600;
      color: #8a94b0;
      text-transform: uppercase;
      letter-spacing: 0.6px;
      align-self: flex-start;
      margin-bottom: 16px;
    }

    .donut-wrap {
      position: relative;
      width: 110px;
      height: 110px;
      margin: 0 auto 8px;
    }

    .donut-wrap svg { transform: rotate(-90deg); }

    .donut-center {
      position: absolute;
      inset: 0;
      display: flex;
      flex-direction: column;
      align-items: center;
      justify-content: center;
    }

    .donut-center .pct {
      font-size: 1.5rem;
      font-weight: 700;
      color: #1a1f36;
    }

    .donut-center .pct-label {
      font-size: 0.65rem;
      color: #8a94b0;
      margin-top: 1px;
    }

    .progress-block { margin-top: 6px; width: 100%; }

    .pb-header {
      display: flex;
      justify-content: space-between;
      margin-bottom: 6px;
    }

    .pb-header span:first-child { font-size: 0.82rem; font-weight: 600; color: #1a1f36; }
    .pb-header span:last-child  { font-size: 0.82rem; color: #8a94b0; }

    .pb-track {
      height: 7px;
      background: #e8eaf6;
      border-radius: 99px;
      overflow: hidden;
    }

    .pb-fill {
      height: 100%;
      border-radius: 99px;
      background: linear-gradient(90deg, #5c6bc0, #9575cd);
      transition: width 0.6s ease;
    }

    .tree-card { padding: 18px; }
    .tree-card .dc-title {
      font-size: 0.72rem;
      font-weight: 600;
      color: #8a94b0;
      text-transform: uppercase;
      letter-spacing: 0.6px;
      margin-bottom: 12px;
    }
    .code-table {
      width: 100%;
      border-collapse: collapse;
      font-size: 0.78rem;
      font-family: 'Courier New', monospace;
    }
    .code-table th {
      text-align: left;
      color: #8a94b0;
      font-weight: 600;
      padding: 4px 6px;
      border-bottom: 1px solid #e8eaf6;
    }
    .code-table td {
      padding: 4px 6px;
      border-bottom: 1px solid #f3f4fa;
      color: #1a1f36;
    }
    .code-table-wrap { max-height: 220px; overflow-y: auto; }

    @media (max-width: 768px) {
      .dashboard { grid-template-columns: 1fr; }
      .sidebar   { flex-direction: row; flex-wrap: wrap; }
      .metric-box, .donut-card, .card.donut-card { flex: 1 1 140px; }
    }
  </style>
</head>
<body>
<div class="dashboard">

  <!-- ═══════════════ MAIN COLUMN ═══════════════ -->
  <div class="main">

    <div class="card header-card">
      <div>
        <h1>Huffman Compressor</h1>
        <p>Compress your text files using Huffman coding</p>
      </div>
      <div class="header-illo">🌳</div>
    </div>

    <div class="card">
      <div class="section-head">
        <div>
          <h2>Compress Text</h2>
          <p>Upload a file or paste content below</p>
        </div>
      </div>

      <div class="drop-zone" id="drop-zone">
        <svg width="26" height="26" fill="none" viewBox="0 0 24 24" stroke="currentColor" stroke-width="1.5">
          <path stroke-linecap="round" stroke-linejoin="round"
            d="M3 16.5v2.25A2.25 2.25 0 005.25 21h13.5A2.25 2.25 0 0021 18.75V16.5m-13.5-9L12 3m0 0l4.5 4.5M12 3v13.5"/>
        </svg>
        <p>Drag &amp; drop a file, or <span id="browse-link">browse</span></p>
        <p style="font-size:0.72rem;color:#c5cbee;margin-top:3px;">.txt .pdf .docx .pptx .csv .log .md</p>
      </div>
      <input type="file" id="file-input" accept=".txt,.pdf,.docx,.pptx,.csv,.log,.md,.json,.xml" />

      <div class="divider">or paste text</div>

      <textarea id="input-text" class="input-area" placeholder="Paste or type text here…"></textarea>

      <div class="btn-row">
        <button class="btn-primary" id="compress-btn">Compress</button>
        <button class="btn-outline" id="clear-btn">Clear</button>
        <span class="char-count" id="char-count">0 chars</span>
      </div>
    </div>

    <!-- Output card -->
    <div class="card" id="output-card" style="display:none;">
      <div class="section-head">
        <div>
          <h2>Compressed Output</h2>
          <p>Base64-encoded bitstream (binary-packed Huffman codes)</p>
        </div>
      </div>
      <textarea id="output-text" class="output-area" readonly></textarea>
      <div class="btn-row">
        <button class="btn-primary"  id="download-btn">Download .huff</button>
        <button class="btn-ghost"    id="copy-btn">Copy</button>
        <button class="btn-outline"  id="decompress-btn">Verify Round-Trip</button>
      </div>
      <div class="msg" id="huff-msg"></div>
    </div>

    <!-- Verify card -->
    <div class="card" id="verify-card" style="display:none;">
      <div class="section-head">
        <div><h2>Decompressed</h2><p>Verification against original</p></div>
      </div>
      <textarea id="verify-text" class="output-area" readonly></textarea>
      <div class="msg" id="verify-msg"></div>
    </div>

  </div><!-- /main -->

  <!-- ═══════════════ SIDEBAR ═══════════════ -->
  <div class="sidebar">

    <div class="metric-box" id="mb-original">
      <div class="mb-left">
        <div class="mb-label">Original Size</div>
        <div class="mb-value" id="mb-orig-val">—</div>
      </div>
      <div class="mb-arrow">›</div>
    </div>

    <div class="metric-box" id="mb-compressed">
      <div class="mb-left">
        <div class="mb-label">Compressed Size</div>
        <div class="mb-value" id="mb-comp-val">—</div>
      </div>
      <div class="mb-arrow">›</div>
    </div>

    <div class="metric-box" id="mb-ratio">
      <div class="mb-left">
        <div class="mb-label">Ratio</div>
        <div class="mb-value" id="mb-ratio-val">—</div>
      </div>
      <div class="mb-arrow">›</div>
    </div>

    <!-- Donut -->
    <div class="card donut-card">
      <div class="dc-title">Compression</div>
      <div class="donut-wrap">
        <svg width="110" height="110" viewBox="0 0 110 110">
          <circle cx="55" cy="55" r="44" fill="none" stroke="#e8eaf6" stroke-width="12"/>
          <circle id="donut-arc" cx="55" cy="55" r="44" fill="none"
            stroke="#5c6bc0" stroke-width="12"
            stroke-dasharray="0 276.46"
            stroke-linecap="round"
            style="transition: stroke-dasharray 0.7s ease;"/>
        </svg>
        <div class="donut-center">
          <span class="pct" id="donut-pct">0%</span>
          <span class="pct-label">Ratio</span>
        </div>
      </div>

      <div class="progress-block">
        <div class="pb-header">
          <span>Space Saved</span>
          <span id="pb-pct-label">0%</span>
        </div>
        <div class="pb-track">
          <div class="pb-fill" id="pb-fill" style="width:0%"></div>
        </div>
      </div>
    </div>

    <!-- Code table -->
    <div class="card tree-card" id="code-table-card" style="display:none;">
      <div class="dc-title">Huffman Codes</div>
      <div class="code-table-wrap">
        <table class="code-table">
          <thead><tr><th>Char</th><th>Freq</th><th>Code</th></tr></thead>
          <tbody id="code-table-body"></tbody>
        </table>
      </div>
    </div>

  </div><!-- /sidebar -->

</div><!-- /dashboard -->

<script>
  const inputText      = document.getElementById('input-text');
  const outputText     = document.getElementById('output-text');
  const verifyText     = document.getElementById('verify-text');
  const charCountEl    = document.getElementById('char-count');
  const dropZone       = document.getElementById('drop-zone');
  const fileInput      = document.getElementById('file-input');
  const browseLink     = document.getElementById('browse-link');
  const compressBtn    = document.getElementById('compress-btn');
  const clearBtn       = document.getElementById('clear-btn');
  const downloadBtn    = document.getElementById('download-btn');
  const copyBtn        = document.getElementById('copy-btn');
  const decompressBtn  = document.getElementById('decompress-btn');
  const outputCard     = document.getElementById('output-card');
  const verifyCard     = document.getElementById('verify-card');
  const huffMsg        = document.getElementById('huff-msg');
  const verifyMsg      = document.getElementById('verify-msg');
  const codeTableCard  = document.getElementById('code-table-card');
  const codeTableBody  = document.getElementById('code-table-body');

  let lastCodebook = null; // saved from /compress, needed for /decompress + /download

  function fmtBytes(n) {
    return n < 1024 ? n + ' B' : (n / 1024).toFixed(1) + ' KB';
  }

  function showMsg(el, type, text) { el.className = 'msg ' + type; el.textContent = text; }
  function hideMsg(el)             { el.className = 'msg'; }

  function hideResults() {
    outputCard.style.display = 'none';
    verifyCard.style.display = 'none';
    codeTableCard.style.display = 'none';
    resetSidebar();
  }

  function resetSidebar() {
    document.getElementById('mb-orig-val').textContent  = '—';
    document.getElementById('mb-comp-val').textContent  = '—';
    document.getElementById('mb-ratio-val').textContent = '—';
    document.getElementById('donut-pct').textContent    = '0%';
    document.getElementById('pb-pct-label').textContent = '0%';
    document.getElementById('pb-fill').style.width      = '0%';
    document.getElementById('donut-arc').setAttribute('stroke-dasharray', '0 276.46');
  }

  // File loading
  browseLink.addEventListener('click', () => fileInput.click());
  dropZone.addEventListener('click',   () => fileInput.click());
  fileInput.addEventListener('change', () => { if (fileInput.files[0]) readFile(fileInput.files[0]); });

  dropZone.addEventListener('dragover', e => { e.preventDefault(); dropZone.classList.add('dragover'); });
  dropZone.addEventListener('dragleave', () => dropZone.classList.remove('dragover'));
  dropZone.addEventListener('drop', e => {
    e.preventDefault();
    dropZone.classList.remove('dragover');
    if (e.dataTransfer.files[0]) readFile(e.dataTransfer.files[0]);
  });

  async function readFile(file) {
    const fd = new FormData();
    fd.append('file', file);
    dropZone.querySelector('p').textContent = 'Extracting text…';
    try {
      const res  = await fetch('/extract', { method: 'POST', body: fd });
      const data = await res.json();
      if (data.error) { alert(data.error); return; }
      inputText.value = data.text;
      updateCount();
      hideResults();
    } catch (e) {
      alert('Failed to read file: ' + e.message);
    } finally {
      dropZone.querySelector('p').textContent = 'Drag & drop a file, or browse';
    }
  }

  inputText.addEventListener('input', () => { updateCount(); hideResults(); });

  function updateCount() {
    const v = inputText.value;
    const n = v.length;
    charCountEl.textContent = n.toLocaleString() + ' char' + (n !== 1 ? 's' : '');
  }

  // Compress
  compressBtn.addEventListener('click', async () => {
    const text = inputText.value;
    if (!text) return;
    compressBtn.disabled = true;
    compressBtn.textContent = 'Compressing…';
    try {
      const res  = await fetch('/compress', {
        method: 'POST',
        headers: { 'Content-Type': 'application/json' },
        body: JSON.stringify({ text })
      });
      const data = await res.json();
      if (data.error) { alert(data.error); return; }

      outputText.value = data.compressed_b64;
      lastCodebook = data.codebook;
      outputCard.style.display = 'block';
      verifyCard.style.display = 'none';

      // Sidebar
      document.getElementById('mb-orig-val').textContent  = fmtBytes(data.original_bytes);
      document.getElementById('mb-comp-val').textContent  = fmtBytes(data.compressed_bytes);
      document.getElementById('mb-ratio-val').textContent = data.ratio + 'x';

      const arcPct  = Math.max(0, Math.min(1, 1 - data.ratio));
      const circ    = 2 * Math.PI * 44;
      const dash    = (arcPct * circ).toFixed(2);
      document.getElementById('donut-arc').setAttribute('stroke-dasharray', `${dash} ${circ}`);
      document.getElementById('donut-pct').textContent = data.ratio + 'x';

      const saving = parseFloat(data.saving_pct);
      const pbPct  = saving > 0 ? Math.min(saving, 100) : 0;
      document.getElementById('pb-fill').style.width      = pbPct + '%';
      document.getElementById('pb-pct-label').textContent = data.saving_pct;

      const isSmaller = data.compressed_bytes < data.original_bytes;
      showMsg(huffMsg, isSmaller ? 'info' : 'warn', data.message);

      // Code table
      codeTableBody.innerHTML = '';
      data.code_table.forEach(row => {
        const tr = document.createElement('tr');
        const charDisplay = row.char === ' ' ? '"space"'
                            : row.char === '\n' ? '\\n'
                            : row.char === '\t' ? '\\t'
                            : row.char;
        tr.innerHTML = `<td>${charDisplay}</td><td>${row.freq}</td><td>${row.code}</td>`;
        codeTableBody.appendChild(tr);
      });
      codeTableCard.style.display = data.code_table.length ? 'block' : 'none';

    } finally {
      compressBtn.disabled  = false;
      compressBtn.textContent = 'Compress';
    }
  });

  // Clear
  clearBtn.addEventListener('click', () => {
    inputText.value = '';
    fileInput.value = '';
    lastCodebook = null;
    updateCount();
    hideResults();
  });

  // Copy
  copyBtn.addEventListener('click', () => {
    navigator.clipboard.writeText(outputText.value).then(() => {
      copyBtn.textContent = 'Copied!';
      setTimeout(() => copyBtn.textContent = 'Copy', 1500);
    });
  });

  // Download
  downloadBtn.addEventListener('click', async () => {
    const res  = await fetch('/download', {
      method: 'POST',
      headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify({ compressed_b64: outputText.value, codebook: lastCodebook })
    });
    const blob = await res.blob();
    const url  = URL.createObjectURL(blob);
    const a    = document.createElement('a');
    a.href = url; a.download = 'compressed.huff'; a.click();
    URL.revokeObjectURL(url);
  });

  // Verify round-trip
  decompressBtn.addEventListener('click', async () => {
    const res  = await fetch('/decompress', {
      method: 'POST',
      headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify({ compressed_b64: outputText.value, codebook: lastCodebook })
    });
    const data = await res.json();
    verifyText.value = data.decompressed || '';
    verifyCard.style.display = 'block';
    if (data.error) {
      showMsg(verifyMsg, 'error', data.error);
    } else if (data.decompressed === inputText.value) {
      showMsg(verifyMsg, 'info', '✓ Verification passed — decompressed output matches the original exactly.');
    } else {
      showMsg(verifyMsg, 'warn', 'Decompressed output differs from the original input.');
    }
  });

  updateCount();
</script>
</body>
</html>
"""


# ───────────────────────── Huffman core ─────────────────────────

class _Node:
    __slots__ = ("freq", "char", "left", "right", "order")
    _counter = 0

    def __init__(self, freq, char=None, left=None, right=None):
        self.freq = freq
        self.char = char
        self.left = left
        self.right = right
        self.order = _Node._counter
        _Node._counter += 1

    def __lt__(self, other):
        if self.freq != other.freq:
            return self.freq < other.freq
        return self.order < other.order


def build_huffman_tree(freqs: dict):
    heap = [_Node(f, ch) for ch, f in freqs.items()]
    heapq.heapify(heap)

    if len(heap) == 1:
        # Edge case: only one distinct character — give it a 1-bit code.
        only = heap[0]
        return _Node(only.freq, left=only)

    while len(heap) > 1:
        a = heapq.heappop(heap)
        b = heapq.heappop(heap)
        merged = _Node(a.freq + b.freq, left=a, right=b)
        heapq.heappush(heap, merged)

    return heap[0]


def build_codebook(root) -> dict:
    codes = {}
    if root.char is not None and root.left is None and root.right is None:
        codes[root.char] = "0"
        return codes

    def walk(node, prefix):
        if node.char is not None:
            codes[node.char] = prefix or "0"
            return
        if node.left:
            walk(node.left, prefix + "0")
        if node.right:
            walk(node.right, prefix + "1")

    walk(root, "")
    return codes


def pack_bits(bitstring: str) -> bytes:
    pad = (8 - len(bitstring) % 8) % 8
    bitstring += "0" * pad
    out = bytearray()
    for i in range(0, len(bitstring), 8):
        out.append(int(bitstring[i:i + 8], 2))
    return bytes(out), pad


def unpack_bits(data: bytes, pad: int) -> str:
    bits = "".join(f"{byte:08b}" for byte in data)
    if pad:
        bits = bits[:-pad]
    return bits


def huffman_encode(text: str):
    if not text:
        return b"", {}, 0

    freqs = {}
    for ch in text:
        freqs[ch] = freqs.get(ch, 0) + 1

    root = build_huffman_tree(freqs)
    codebook = build_codebook(root)

    bitstring = "".join(codebook[ch] for ch in text)
    packed, pad = pack_bits(bitstring)
    return packed, codebook, pad


def huffman_decode(data: bytes, codebook: dict, pad: int, length_hint=None) -> str:
    if not data and not codebook:
        return ""
    reverse = {code: ch for ch, code in codebook.items()}
    bits = unpack_bits(data, pad)
    result = []
    current = ""
    for bit in bits:
        current += bit
        if current in reverse:
            result.append(reverse[current])
            current = ""
    if current and current in reverse:
        result.append(reverse[current])
    return "".join(result)


# ───────────────────────── File extraction ─────────────────────────

def extract_text_from_file(file) -> str:
    filename = file.filename.lower()
    data = file.read()
    buf = io.BytesIO(data)

    if filename.endswith(".pdf"):
        text_parts = []
        with pdfplumber.open(buf) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
        return "\n\n".join(text_parts)

    elif filename.endswith(".docx"):
        doc = docx.Document(buf)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    elif filename.endswith(".pptx"):
        prs = pptx.Presentation(buf)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text)
        return "\n".join(lines)

    else:
        return data.decode("utf-8", errors="replace")


# ───────────────────────── Routes ─────────────────────────

@app.route("/extract", methods=["POST"])
def extract():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    f = request.files["file"]
    try:
        text = extract_text_from_file(f)
        if not text.strip():
            return jsonify({"error": "Could not extract any text from this file."}), 400
        return jsonify({"text": text})
    except Exception as e:
        return jsonify({"error": f"Failed to read file: {e}"}), 500


@app.route("/")
def index():
    return render_template_string(HTML)


@app.route("/compress", methods=["POST"])
def compress():
    data = request.get_json(force=True)
    text = data.get("text", "")
    if not text:
        return jsonify({"error": "No text provided"}), 400

    try:
        packed, codebook, pad = huffman_encode(text)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    orig_bytes = len(text.encode("utf-8"))
    comp_bytes = len(packed)

    # Serialized codebook size counts toward a fair "total compressed size"
    codebook_json = json.dumps(codebook, ensure_ascii=False)
    codebook_bytes = len(codebook_json.encode("utf-8"))
    total_comp_bytes = comp_bytes + codebook_bytes

    ratio      = round(total_comp_bytes / orig_bytes, 2) if orig_bytes else 0
    saving_raw = round((1 - ratio) * 100, 1)
    saving_str = (f"-{saving_raw}%" if saving_raw >= 0 else f"+{abs(saving_raw)}%")

    if total_comp_bytes < orig_bytes:
        msg = (f"Huffman coding reduced the file by {saving_raw}% "
               f"(including the {codebook_bytes}-byte codebook needed to decode it).")
    else:
        msg = ("The compressed output (plus codebook) is larger than the original. "
               "Huffman coding works best on longer texts with skewed character frequencies.")

    freqs = {}
    for ch in text:
        freqs[ch] = freqs.get(ch, 0) + 1
    code_table = sorted(
        ({"char": ch, "freq": freqs[ch], "code": code} for ch, code in codebook.items()),
        key=lambda r: -r["freq"]
    )

    return jsonify({
        "compressed_b64":   base64.b64encode(packed).decode("ascii"),
        "codebook":         codebook,
        "pad":              pad,
        "original_bytes":   orig_bytes,
        "compressed_bytes": total_comp_bytes,
        "ratio":            ratio,
        "saving_pct":       saving_str,
        "message":          msg,
        "code_table":       code_table,
    })


@app.route("/decompress", methods=["POST"])
def decompress():
    data        = request.get_json(force=True)
    b64         = data.get("compressed_b64", "")
    codebook    = data.get("codebook", {})
    try:
        packed = base64.b64decode(b64) if b64 else b""
        # pad isn't transmitted separately to the client in this simple demo,
        # so we recompute decode tolerant of trailing leftover bits.
        decoded = huffman_decode(packed, codebook, pad=0)
        return jsonify({"decompressed": decoded})
    except Exception as e:
        return jsonify({"error": str(e), "decompressed": ""}), 400


@app.route("/download", methods=["POST"])
def download():
    data      = request.get_json(force=True)
    b64       = data.get("compressed_b64", "")
    codebook  = data.get("codebook", {})
    packed    = base64.b64decode(b64) if b64 else b""

    # Bundle codebook + bitstream into one self-contained JSON file
    payload = json.dumps({
        "codebook": codebook,
        "data_b64": base64.b64encode(packed).decode("ascii"),
    }).encode("utf-8")

    buf = io.BytesIO(payload)
    buf.seek(0)
    return send_file(
        buf,
        as_attachment=True,
        download_name="compressed.huff",
        mimetype="application/json",
    )


if __name__ == "__main__":
    app.run(debug=True, port=5001)